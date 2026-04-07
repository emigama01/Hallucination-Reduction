"""
Generate updated PowerPoint presentation with experimental results.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
import json
import os

RESULTS_DIR = os.path.join(os.path.dirname(__file__), "..", "results")
OUTPUT_DIR = os.path.dirname(os.path.dirname(__file__))


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(content_lines):
        if i == 0:
            tf.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0 if not line.startswith("  ") else 1
    return slide


def add_blank_slide_with_title(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    # Add title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    p.alignment = PP_ALIGN.LEFT
    return slide


def add_image_slide(prs, title, image_path, top=1.5, left=1.0, width=8.0):
    slide = add_blank_slide_with_title(prs, title)
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width))
    return slide


def generate_ppt():
    with open(os.path.join(RESULTS_DIR, "summary.json")) as f:
        summary = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title ----
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hallucination Reduction in\nLarge Language Models"
    slide.placeholders[1].text = (
        "A Comparative Study of Prompt Engineering and RAG\n\n"
        "Emilio Garcia Martinez | AI700-001\n"
        "Professor: Reda Nacif Elalaoui"
    )

    # ---- Slide 2: Problem Statement ----
    add_content_slide(prs, "Problem Statement", [
        "LLMs generate fluent but factually incorrect content (hallucinations)",
        "Critical risk in high-stakes domains: healthcare, legal, education",
        "Need systematic evaluation of reduction techniques",
        "Research Question: Which techniques most effectively reduce hallucinations?",
        "  Focus: Prompt engineering vs. Retrieval-Augmented Generation (RAG)",
    ])

    # ---- Slide 3: Approach Overview ----
    add_content_slide(prs, "Three-Phase Experimental Approach", [
        "Phase 1: Baseline Measurement",
        "  Vanilla LLaMA 3.2 (3B) on 30 TruthfulQA-inspired questions",
        "Phase 2: Prompt-Based Reduction",
        "  Chain-of-Thought (CoT), Few-Shot, Self-Consistency Decoding",
        "Phase 3: RAG-Based Reduction",
        "  ChromaDB knowledge base + optimized prompting pipeline",
        "Evaluation: ROUGE-L, Keyword Recall, Factual Accuracy Score",
    ])

    # ---- Slide 4: Dataset ----
    add_content_slide(prs, "Dataset & Evaluation", [
        "30 curated TruthfulQA-inspired questions across 5 categories:",
        "  Health (7) | Science (8) | History (5) | Geography (4) | Misconceptions (6)",
        "Each question has verified correct answer + known incorrect alternatives",
        "Metrics:",
        "  Hallucination Detection: keyword overlap + ROUGE-L comparison",
        "  Factual Accuracy: composite (0.6 keyword recall + 0.4 ROUGE-L)",
        "  ROUGE-L: lexical overlap with ground truth",
    ])

    # ---- Slide 5: Phase 1 Results ----
    baseline = summary['baseline']
    add_content_slide(prs, "Phase 1: Baseline Results", [
        f"Model: LLaMA 3.2 (3B parameters) via Ollama",
        f"Hallucination Rate: {baseline['hallucination_rate']:.1%}",
        f"Average Factual Accuracy: {baseline['avg_factual_accuracy']:.3f}",
        f"Average ROUGE-L: {baseline['avg_rouge_l']:.3f}",
        "Highest hallucination categories:",
        f"  History: {baseline['by_category']['History']['hallucination_rate']:.0%} hallucination rate",
        f"  Misconceptions: {baseline['by_category']['Misconceptions']['hallucination_rate']:.0%} hallucination rate",
        f"  Health: {baseline['by_category']['Health']['hallucination_rate']:.0%} hallucination rate (lowest)",
    ])

    # ---- Slide 6: Phase 2 Results ----
    add_content_slide(prs, "Phase 2: Prompt Engineering Results", [
        f"Chain-of-Thought: Accuracy {summary['cot']['avg_factual_accuracy']:.3f} | ROUGE-L {summary['cot']['avg_rouge_l']:.3f}",
        "  Verbose outputs reduce ROUGE-L but factual content is often correct",
        f"Few-Shot Prompting: Accuracy {summary['few_shot']['avg_factual_accuracy']:.3f} | ROUGE-L {summary['few_shot']['avg_rouge_l']:.3f}",
        "  Best prompt-only method for factual accuracy (+18.5% over baseline)",
        f"Self-Consistency: Accuracy {summary['self_consistency']['avg_factual_accuracy']:.3f} | ROUGE-L {summary['self_consistency']['avg_rouge_l']:.3f}",
        "  Best hallucination rate reduction among prompt methods",
        "Key Finding: Prompt engineering alone has limited effect on 3B model",
    ])

    # ---- Slide 7: Phase 3 Results ----
    rag_imp = ((summary['rag_optimized']['avg_factual_accuracy'] - baseline['avg_factual_accuracy'])
               / baseline['avg_factual_accuracy'] * 100)
    add_content_slide(prs, "Phase 3: RAG Pipeline Results", [
        f"RAG Pipeline: Accuracy {summary['rag']['avg_factual_accuracy']:.3f} | ROUGE-L {summary['rag']['avg_rouge_l']:.3f}",
        "  Highest ROUGE-L score across all methods",
        f"RAG + Optimized: Accuracy {summary['rag_optimized']['avg_factual_accuracy']:.3f} | ROUGE-L {summary['rag_optimized']['avg_rouge_l']:.3f}",
        f"  Highest factual accuracy: +{rag_imp:.1f}% improvement over baseline",
        "RAG Knowledge Base: 4 documents, 27 chunks, ChromaDB + MiniLM-L6-v2",
        "Key Finding: Grounding in external knowledge is most effective approach",
    ])

    # ---- Slide 8: Hallucination Rates Chart ----
    img_path = os.path.join(RESULTS_DIR, "hallucination_rates.png")
    add_image_slide(prs, "Hallucination Rates Across Methods", img_path, top=1.3, left=2.5, width=8.0)

    # ---- Slide 9: Accuracy Comparison Chart ----
    img_path = os.path.join(RESULTS_DIR, "accuracy_comparison.png")
    add_image_slide(prs, "Factual Accuracy & ROUGE-L Comparison", img_path, top=1.3, left=2.5, width=8.0)

    # ---- Slide 10: Category Heatmap ----
    img_path = os.path.join(RESULTS_DIR, "category_heatmap.png")
    add_image_slide(prs, "Category-Level Hallucination Analysis", img_path, top=1.3, left=2.5, width=8.0)

    # ---- Slide 11: Key Findings ----
    add_content_slide(prs, "Key Findings", [
        f"1. RAG + Optimized Prompting achieves {rag_imp:.1f}% factual accuracy improvement",
        "2. Prompt engineering alone has limited effect on smaller (3B) models",
        "3. Self-consistency is the most practical prompt-only improvement",
        "4. ROUGE-L underestimates structured prompting effectiveness",
        "5. RAG is most effective for factual recall (Geography, Health)",
        "6. Misconceptions and History are hardest categories to address",
        "7. External knowledge grounding > parametric knowledge for accuracy",
    ])

    # ---- Slide 12: Limitations & Future Work ----
    add_content_slide(prs, "Limitations & Future Work", [
        "Limitations:",
        "  Small dataset (30 questions), single model size (3B)",
        "  Purpose-built knowledge base (optimistic RAG scenario)",
        "  Automated evaluation only (no human review yet)",
        "Future Work:",
        "  Evaluate on larger models (7B, 13B, 70B parameters)",
        "  Implement hallucination-aware fine-tuning (RAG-HAT / DPO)",
        "  Full TruthfulQA benchmark evaluation",
        "  NLI-based faithfulness scoring",
    ])

    # ---- Slide 13: References ----
    add_content_slide(prs, "References", [
        "[1] Dang et al., Survey of Hallucinations in LLMs, Frontiers in AI, 2025",
        "[2] Bang et al., HalluLens: LLM Hallucination Benchmark, ACL 2025",
        "[3] Wang et al., Self-Consistency in Chain-of-Thought, ICLR 2023",
        "[4] Song et al., RAG-HAT Pipeline, EMNLP 2024",
        "[5] Lin et al., TruthfulQA, ACL 2022",
        "[6] Lewis et al., Retrieval-Augmented Generation, NeurIPS 2020",
        "[7] Wei et al., Chain-of-Thought Prompting, NeurIPS 2022",
    ])

    # ---- Slide 14: Thank You ----
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You"
    slide.placeholders[1].text = (
        "Questions?\n\n"
        "Emilio Garcia Martinez\n"
        "AI700-001 | ID: 100783029"
    )

    output_path = os.path.join(OUTPUT_DIR, "AI700-001-Progress-PPT_Hallucination-LLM.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_ppt()
