"""
Update the original college-formatted PPT with actual experimental results.
Preserves college format/theme, updates content and adds result slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import os
import copy
from lxml import etree

RESULTS_DIR = os.path.join(os.path.dirname(__file__), "..", "results")
INPUT_PPT = os.path.join(os.path.dirname(__file__), "..", "AI700-001-Proposal-PPT_Hallucination-LLM.pptx")
OUTPUT_PPT = os.path.join(os.path.dirname(__file__), "..", "AI700-001-Progress-PPT_Hallucination-LLM.pptx")


def load_results():
    with open(os.path.join(RESULTS_DIR, "summary.json")) as f:
        return json.load(f)


def clear_placeholder(placeholder):
    """Clear all text from a placeholder while keeping formatting of first paragraph."""
    tf = placeholder.text_frame
    # Keep the first paragraph's format reference
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    # Clear first paragraph text
    tf.paragraphs[0].clear()


def set_bullet_text(placeholder, lines, font_size=14, bold_indices=None, sub_indices=None):
    """Set bullet-point text in a placeholder, preserving format from template."""
    if bold_indices is None:
        bold_indices = set()
    if sub_indices is None:
        sub_indices = set()

    tf = placeholder.text_frame
    tf.word_wrap = True

    # Clear existing content
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    tf.paragraphs[0].clear()

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullet indentation
        if i in sub_indices:
            p.level = 1
            actual_text = line.lstrip()
        else:
            p.level = 0
            actual_text = line

        run = p.add_run()
        run.text = actual_text
        run.font.size = Pt(font_size)
        run.font.name = 'Calibri'

        if i in bold_indices:
            run.font.bold = True
        else:
            run.font.bold = False


def add_slide_with_format(prs, layout_idx, title_text, ref_slide_idx=2):
    """Add a new slide using the specified layout and matching the title format."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    return slide


def set_title_box(slide, title_text):
    """Find the title rectangle and set its text, or create one."""
    for shape in slide.shapes:
        if shape.name == 'Rectangle 2' and shape.has_text_frame:
            tf = shape.text_frame
            if tf.paragraphs:
                tf.paragraphs[0].clear()
                run = tf.paragraphs[0].add_run()
                run.text = title_text
                # Copy formatting from the text
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.name = 'Calibri'
            return True
    return False


def create_title_box(slide, title_text, ref_shape):
    """Create a title box matching the reference shape's position/size."""
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(
        ref_shape.left, ref_shape.top,
        ref_shape.width, ref_shape.height
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = 'Calibri'
    return txBox


def update_slide_number(slide, num):
    """Update slide number placeholder."""
    for shape in slide.shapes:
        if 'Slide Number' in shape.name and shape.has_text_frame:
            shape.text_frame.paragraphs[0].clear()
            run = shape.text_frame.paragraphs[0].add_run()
            run.text = str(num)
            run.font.size = Pt(12)


def main():
    summary = load_results()
    prs = Presentation(INPUT_PPT)

    slides = list(prs.slides)

    # ============================================================
    # Update Slide 7 (index 6): Methodology - update to reflect actual work
    # ============================================================
    slide7 = slides[6]
    for shape in slide7.shapes:
        if shape.name == 'Text Placeholder 3' and shape.has_text_frame:
            set_bullet_text(shape, [
                "Phase 1 - Baseline Measurement:",
                "Measured hallucination rates using LLaMA 3.2 (3B) on 30 TruthfulQA-inspired questions across 5 categories (Health, Science, History, Geography, Misconceptions)",
                "Phase 2 - Prompt-Based Reduction:",
                "Implemented chain-of-thought (CoT), few-shot prompting with verified examples, and self-consistency decoding (3 samples with majority vote)",
                "Phase 3 - RAG-Based Reduction Pipeline:",
                "Built RAG pipeline with ChromaDB + all-MiniLM-L6-v2 embeddings (27 indexed chunks), tested both vanilla RAG and RAG + optimized prompting combining CoT + few-shot + retrieval",
            ], font_size=13, bold_indices={0, 2, 4})
            break

    # ============================================================
    # Update Slide 8 (index 7): Tools - update to reflect actual tools
    # ============================================================
    slide8 = slides[7]
    for shape in slide8.shapes:
        if shape.name == 'Text Placeholder 3' and shape.has_text_frame:
            set_bullet_text(shape, [
                "Language: Python 3.9+",
                "LLM: LLaMA 3.2 (3B parameters) via Ollama (local inference)",
                "RAG Framework: ChromaDB for vector storage + retrieval",
                "Embeddings: all-MiniLM-L6-v2 (sentence-transformers)",
                "Evaluation: ROUGE-L, Keyword Recall, Composite Factual Accuracy",
                "Data Processing: pandas, NumPy",
                "Visualization: matplotlib, seaborn",
                "Environment: Local Apple M1 (8GB RAM), Jupyter-compatible scripts",
            ], font_size=13)
            break

    # ============================================================
    # Update Slide 9 (index 8): Results - replace with actual data
    # ============================================================
    slide9 = slides[8]

    # Remove the old placeholder image
    shapes_to_remove = []
    for shape in slide9.shapes:
        if shape.shape_type == 13:  # PICTURE
            shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Update the content placeholder with actual results
    for shape in slide9.shapes:
        if shape.name == 'Content Placeholder 5' and shape.has_text_frame:
            # Reposition to use full slide area
            shape.top = Emu(1131013)
            shape.height = Emu(5200000)

            baseline_acc = summary['baseline']['avg_factual_accuracy']
            best_acc = summary['rag_optimized']['avg_factual_accuracy']
            improvement = ((best_acc - baseline_acc) / baseline_acc) * 100

            set_bullet_text(shape, [
                "Experimental Results (LLaMA 3.2 3B, 30 TruthfulQA-inspired questions):",
                f"Vanilla Baseline: Factual Accuracy 0.388 | ROUGE-L 0.283 | Halluc. Rate 66.7%",
                f"Chain-of-Thought: Factual Accuracy 0.346 | ROUGE-L 0.102 (verbose outputs dilute ROUGE)",
                f"Few-Shot Prompting: Factual Accuracy 0.460 | ROUGE-L 0.159 (best prompt-only method)",
                f"Self-Consistency: Factual Accuracy 0.448 | ROUGE-L 0.274 (best halluc. rate reduction)",
                f"RAG Pipeline: Factual Accuracy 0.510 | ROUGE-L 0.353 (highest ROUGE-L)",
                f"RAG + Optimized: Factual Accuracy 0.556 | ROUGE-L 0.195 (highest accuracy, +{improvement:.1f}%)",
                "Key Finding: RAG grounding is most effective; prompt-only methods limited on 3B model",
            ], font_size=12, bold_indices={0, 7})
            break

    # Update title
    set_title_box(slide9, "EXPERIMENTAL RESULTS")

    # ============================================================
    # Now we need to add new slides AFTER slide 9 for charts
    # We'll add them after the results slide
    # ============================================================

    # Get reference shapes for formatting new slides
    ref_title_shape = None
    for shape in slides[2].shapes:
        if shape.name == 'Rectangle 2':
            ref_title_shape = shape
            break

    # --- New Slide: Hallucination Rates Chart ---
    new_slide1 = add_slide_with_format(prs, 4, "")  # Blank layout
    # Add title box by copying XML from reference
    txBox1 = new_slide1.shapes.add_textbox(
        Emu(628650), Emu(30063), Emu(6536079), Emu(870891)
    )
    tf1 = txBox1.text_frame
    p1 = tf1.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "HALLUCINATION RATES COMPARISON"
    run1.font.size = Pt(28)
    run1.font.bold = True
    run1.font.name = 'Calibri'

    # Add chart image
    img_path = os.path.join(RESULTS_DIR, "hallucination_rates.png")
    if os.path.exists(img_path):
        new_slide1.shapes.add_picture(img_path, Emu(1200000), Emu(1100000), Emu(6800000))

    # --- New Slide: Accuracy Comparison Chart ---
    new_slide2 = add_slide_with_format(prs, 4, "")
    txBox2 = new_slide2.shapes.add_textbox(
        Emu(628650), Emu(30063), Emu(6536079), Emu(870891)
    )
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "ACCURACY & ROUGE-L COMPARISON"
    run2.font.size = Pt(28)
    run2.font.bold = True
    run2.font.name = 'Calibri'

    img_path2 = os.path.join(RESULTS_DIR, "accuracy_comparison.png")
    if os.path.exists(img_path2):
        new_slide2.shapes.add_picture(img_path2, Emu(1200000), Emu(1100000), Emu(6800000))

    # --- New Slide: Category Heatmap ---
    new_slide3 = add_slide_with_format(prs, 4, "")
    txBox3 = new_slide3.shapes.add_textbox(
        Emu(628650), Emu(30063), Emu(6536079), Emu(870891)
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "CATEGORY-LEVEL ANALYSIS"
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.name = 'Calibri'

    img_path3 = os.path.join(RESULTS_DIR, "category_heatmap.png")
    if os.path.exists(img_path3):
        new_slide3.shapes.add_picture(img_path3, Emu(1200000), Emu(1100000), Emu(6800000))

    # --- New Slide: Key Findings ---
    new_slide4 = add_slide_with_format(prs, 3, "")  # Title and Content layout
    set_title_box(new_slide4, "KEY FINDINGS")

    baseline_acc = summary['baseline']['avg_factual_accuracy']
    best_acc = summary['rag_optimized']['avg_factual_accuracy']
    improvement = ((best_acc - baseline_acc) / baseline_acc) * 100

    for shape in new_slide4.shapes:
        if 'Content Placeholder' in shape.name and shape.has_text_frame:
            set_bullet_text(shape, [
                f"RAG + Optimized Prompting achieves {improvement:.1f}% factual accuracy improvement over baseline",
                "Prompt engineering alone has limited effect on smaller (3B parameter) models",
                "Self-consistency decoding is the most practical prompt-only improvement (3x cost)",
                "ROUGE-L alone underestimates effectiveness of structured prompting (CoT/RAG+Opt)",
                "RAG is most effective for factual recall categories (Geography, Health)",
                "History and Misconceptions are the hardest categories to address with current methods",
                "Grounding outputs in external knowledge > relying on parametric model knowledge",
            ], font_size=13)
            break

    # ============================================================
    # Update Slide 10 (now further back): Conclusion
    # ============================================================
    # Find the conclusion slide (originally index 9)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == 'Rectangle 2' and shape.has_text_frame:
                if 'CONCLUSION' in shape.text_frame.text:
                    # Found conclusion slide - update content
                    for s in slide.shapes:
                        if s.name == 'Content Placeholder 2' and s.has_text_frame:
                            set_bullet_text(s, [
                                "Hallucination reduction is critical for trustworthy LLM deployment",
                                f"RAG + Optimized Prompting achieved {improvement:.1f}% factual accuracy improvement",
                                "Few-shot and self-consistency are best prompt-only techniques",
                                "Chain-of-thought less effective on smaller models (capacity-dependent)",
                                "No single technique eliminates hallucinations - combining strategies is key",
                                "Future Work:",
                                "  Evaluate on larger models (7B, 13B, 70B parameters)",
                                "  Implement hallucination-aware fine-tuning (RAG-HAT / DPO)",
                                "  Full TruthfulQA benchmark + NLI-based faithfulness scoring",
                                "  Domain-specific reduction pipelines (medical, legal)",
                            ], font_size=13, bold_indices={5}, sub_indices={6, 7, 8, 9})
                            break
                    break

    # ============================================================
    # Update References slide
    # ============================================================
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == 'Rectangle 2' and shape.has_text_frame:
                if 'REFERENCES' in shape.text_frame.text:
                    for s in slide.shapes:
                        if s.name == 'Content Placeholder 2' and s.has_text_frame:
                            set_bullet_text(s, [
                                "[1] Dang et al. (2025). Survey of Hallucinations in LLMs. Frontiers in AI, vol. 8.",
                                "[2] Bang et al. (2025). HalluLens: LLM Hallucination Benchmark. ACL 2025.",
                                "[3] Song et al. (2024). RAG-HAT: Hallucination-Aware Tuning Pipeline. EMNLP 2024.",
                                "[4] Wang et al. (2023). Self-Consistency in Chain-of-Thought Reasoning. ICLR 2023.",
                                "[5] Lin et al. (2022). TruthfulQA: Measuring How Models Mimic Falsehoods. ACL 2022.",
                                "[6] Lewis et al. (2020). Retrieval-Augmented Generation for NLP. NeurIPS 2020.",
                                "[7] Wei et al. (2022). Chain-of-Thought Prompting. NeurIPS 2022.",
                            ], font_size=12)
                            break
                    break

    # ============================================================
    # Reorder slides: move new slides (added at end) to after slide 9
    # Original: 1-9, 10(conclusion), 11(refs), 12(questions), NEW1, NEW2, NEW3, NEW4
    # Desired:  1-9, NEW1, NEW2, NEW3, NEW4, 10, 11, 12
    # ============================================================

    # The new slides are at indices 12, 13, 14, 15 (0-indexed)
    # We want them after index 8 (slide 9)

    sldIdLst = prs.part._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst'
    )
    sldIds = list(sldIdLst)

    # Move slides 12,13,14,15 to after position 8
    # Current order: 0,1,2,3,4,5,6,7,8,9,10,11,12,13,14,15
    # Desired:       0,1,2,3,4,5,6,7,8,12,13,14,15,9,10,11

    if len(sldIds) >= 16:
        new_order = list(range(9)) + [12, 13, 14, 15] + [9, 10, 11]
        reordered = [sldIds[i] for i in new_order]

        for sldId in sldIds:
            sldIdLst.remove(sldId)
        for sldId in reordered:
            sldIdLst.append(sldId)

    # Update all slide numbers
    for idx, slide in enumerate(prs.slides):
        update_slide_number(slide, idx + 1)

    # Save
    prs.save(OUTPUT_PPT)
    print(f"Updated PPT saved to: {OUTPUT_PPT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
