"""Fill the final PPT template with project content (12-slide structure)."""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "AI700-001_BK_SPRG2026-PPT-Template-Final.pptx"
OUT = ROOT / "AI700-Final-PPT-Hallucination-LLM.pptx"

STATS = json.loads((ROOT / "results" / "stats_summary.json").read_text())


def set_tf(shape, lines, body_size=18, first_size=None, bold_first=False):
    """Set the text frame of a shape with multiple paragraphs."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # don't reset the bullet level - keep template default
        run = p.add_run()
        run.text = line
        size = first_size if (i == 0 and first_size) else body_size
        run.font.size = Pt(size)
        if i == 0 and bold_first:
            run.font.bold = True


def find_shape(slide, name=None, contains_text=None, placeholder_idx=None):
    for shape in slide.shapes:
        if name and shape.name == name:
            return shape
        if placeholder_idx is not None and shape.is_placeholder and shape.placeholder_format.idx == placeholder_idx:
            return shape
        if contains_text and shape.has_text_frame and contains_text in shape.text_frame.text:
            return shape
    return None


def replace_image(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def main():
    p = Presentation(str(TEMPLATE))
    slides = list(p.slides)

    # ---- Slide 1: Title ----
    s1 = slides[0]
    title_shape = find_shape(s1, name="Rectangle 2")
    if title_shape and title_shape.has_text_frame:
        set_tf(title_shape, ["Hallucination Reduction in Large Language Models",
                             "A Comparative Study of Prompt Engineering and RAG"],
               body_size=24, first_size=28, bold_first=True)
    info_shape = find_shape(s1, name="TextBox 8")
    if info_shape:
        set_tf(info_shape, [
            "By: Emilio Garcia Martinez (ID: 100783029)",
            "Course: AI700-001 | Prof. Reda Nacif Elalaoui",
            "Date: May 2026 | Spring 2026",
        ], body_size=16)
    # Replace the "Picture Representing The Scope Of Your Project" placeholder
    # with the actual project hero chart.
    pic_shape = find_shape(s1, name="TextBox 5")
    if pic_shape and pic_shape.has_text_frame:
        left, top, width, height = pic_shape.left, pic_shape.top, pic_shape.width, pic_shape.height
        sp = pic_shape._element
        sp.getparent().remove(sp)
        hero_img = ROOT / "results" / "accuracy_with_ci.png"
        if hero_img.exists():
            s1.shapes.add_picture(str(hero_img), left, top, width=width, height=height)

    # ---- Slide 2: Outlines ----
    s2 = slides[1]
    outline = find_shape(s2, placeholder_idx=1)
    if outline:
        set_tf(outline, [
            "1. Introduction & Motivation",
            "2. Application & Real-World Importance",
            "3. Literature Review",
            "4. Literature Review Results",
            "5. Methodology — Pipeline & Methods",
            "6. Understanding the Functions (Metrics)",
            "7. Obtained Results & Statistical Analysis",
            "8. Conclusion & Future Work",
            "9. References",
        ], body_size=20)

    # ---- Slide 3: Introduction ----
    s3 = slides[2]
    body = find_shape(s3, placeholder_idx=1)
    if body:
        set_tf(body, [
            "Hallucination: LLMs generating fluent but factually incorrect text.",
            "Persistent failure mode in LLaMA, GPT, Claude — high-stakes risk in medicine, law, education.",
            "Goal: compare six mitigation techniques on LLaMA 3.2 (3B) under one protocol.",
            "Six methods: Baseline, Chain-of-Thought, Few-Shot, Self-Consistency, RAG, RAG + Optimized.",
            "Evaluation: 30-question TruthfulQA subset, automated metrics, bootstrap statistics.",
        ], body_size=18)

    # ---- Slide 4: Application ----
    s4 = slides[3]
    body = find_shape(s4, placeholder_idx=1)
    if body:
        set_tf(body, [
            "Healthcare assistants — wrong dosage / drug-interaction advice can harm patients.",
            "Legal & finance copilots — fabricated case law or numbers create liability.",
            "Educational tutors — myth reinforcement in History / Science misinforms students.",
            "Customer support automation — confidently wrong answers erode trust.",
            "Reducing hallucination is foundational to deploying LLMs in any production setting.",
        ], body_size=18)

    # ---- Slide 5: Literature Review ----
    s5 = slides[4]
    body = find_shape(s5, placeholder_idx=2)
    if body:
        set_tf(body, [
            "Dang et al. (2025): survey & taxonomy of LLM hallucinations [1].",
            "Bang et al. (2025): HalluLens benchmark for systematic evaluation [2].",
            "Wang et al. (2023): Self-Consistency boosts CoT reasoning [3].",
            "Song et al. (2024): RAG-HAT — hallucination-aware tuning for RAG [4].",
            "Lin et al. (2022): TruthfulQA targets imitative falsehoods [5].",
            "Lewis et al. (2020): Retrieval-Augmented Generation framework [7].",
            "Wei et al. (2022): Chain-of-Thought prompting elicits reasoning [8].",
        ], body_size=16)

    # ---- Slide 6: Lit Review Results ----
    s6 = slides[5]
    body = find_shape(s6, placeholder_idx=1)
    if body:
        set_tf(body, [
            "RAG and verification methods consistently outperform pure prompt engineering on factual tasks.",
            "Self-Consistency's value depends on cost tolerance — 3+ inference passes per query.",
            "CoT benefits scale with model size; small models (≤7B) often see marginal or negative gains.",
            "Open evaluation gap: few studies report paired statistical tests on small open-weight models running locally.",
            "This work contributes a fully reproducible local pipeline with bootstrap CIs and error-type analysis.",
        ], body_size=18)

    # ---- Slide 7: Methodology ----
    s7 = slides[6]
    body = find_shape(s7, placeholder_idx=2)
    if body:
        set_tf(body, [
            "Model: LLaMA 3.2 (3B) via Ollama on Apple M1 (8 GB RAM).",
            "Dataset: 30 TruthfulQA-inspired questions across 5 categories.",
            "Six methods: Baseline, CoT, Few-Shot, Self-Consistency, RAG, RAG+Opt.",
            "RAG knowledge base: 27 chunks in ChromaDB; top-3 retrieval.",
            "Evaluation: factual accuracy, ROUGE-L, hallucination rate.",
            "Statistics: 10k bootstrap CIs and paired bootstrap vs baseline.",
        ], body_size=14)
    # Add methodology pipeline image
    pipeline_img = ROOT / "results" / "radar_comparison.png"
    if pipeline_img.exists():
        # Find the picture placeholder (Content Placeholder 2 with no text frame)
        for shape in s7.shapes:
            if shape.name == "Content Placeholder 2" and not shape.has_text_frame:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                # Remove the placeholder
                sp = shape._element
                sp.getparent().remove(sp)
                replace_image(s7, pipeline_img, left, top, width, height)
                break

    # ---- Slide 8: Understanding the Functions (metrics) ----
    s8 = slides[7]
    body = find_shape(s8, placeholder_idx=2)
    if body:
        set_tf(body, [
            "Hallucination Rate: 1 if response overlaps known-incorrect more than correct answer.",
            "ROUGE-L: longest common subsequence overlap between response and ground truth.",
            "Factual Accuracy = 0.6 × keyword recall + 0.4 × ROUGE-L (rewards salient facts under paraphrase).",
            "Bootstrap 95% CI: 10,000 resamples per method.",
            "Paired bootstrap test: per-question Δ vs baseline; two-sided p from sign-flip fraction.",
        ], body_size=16)

    # ---- Slide 9: Obtained Results ----
    s9 = slides[8]
    body = find_shape(s9, placeholder_idx=1)
    if body:
        # Build a results summary
        lines = [
            "Best method: RAG + Optimized (factual accuracy 0.556, +43.3% over baseline).",
            "RAG + Optimized:  0.556 [0.489, 0.624]   (+0.168 vs baseline, p = 0.001)",
            "RAG:                 0.510 [0.417, 0.605]   (+0.121, p = 0.031)",
            "Self-Consistency: 0.448 [0.385, 0.514]   (+0.060, p = 0.034)",
            "Few-Shot:           0.460 [0.394, 0.527]   (+0.072, p = 0.125, n.s.)",
            "Chain-of-Thought: 0.346 [0.292, 0.408]   (-0.042, p = 0.340, n.s.)",
            "Baseline:           0.388 [0.313, 0.468]",
            "Error analysis: residual failures dominated by partial-truth, not full fabrication.",
        ]
        set_tf(body, lines, body_size=14)
    # Insert results chart as picture on the same slide
    chart_img = ROOT / "results" / "accuracy_with_ci.png"
    if chart_img.exists():
        slide_w = p.slide_width
        slide_h = p.slide_height
        img_w = Inches(5.5)
        img_h = Inches(3.1)
        # bottom-right corner
        left = slide_w - img_w - Inches(0.3)
        top = slide_h - img_h - Inches(0.4)
        s9.shapes.add_picture(str(chart_img), left, top, width=img_w, height=img_h)

    # ---- Slide 10: Conclusion & Future Work ----
    s10 = slides[9]
    body = find_shape(s10, placeholder_idx=1)
    if body:
        set_tf(body, [
            "RAG + Optimized prompting is the most reliable hallucination mitigation on LLaMA 3.2 (3B).",
            "RAG and Self-Consistency also produce statistically significant gains.",
            "Chain-of-Thought is not effective at this model size — confirms scale-dependence findings.",
            "Residual errors are mostly partial-truth — model surfaces the topic but fails to commit to the verified fact.",
            "Future work: factual-entailment verification, hallucination-aware fine-tuning [4], and larger-model replication.",
            "Public release: full pipeline + dataset + results on GitHub for reproducibility.",
        ], body_size=16)

    # ---- Slide 11: References ----
    s11 = slides[10]
    body = find_shape(s11, placeholder_idx=1)
    if body:
        set_tf(body, [
            "[1] Dang, Vu, Nguyen. \"Survey on hallucinations in LLMs.\" Frontiers in AI, 2025.",
            "[2] Bang et al. \"HalluLens benchmark.\" ACL 2025.",
            "[3] Wang et al. \"Self-Consistency improves CoT.\" ICLR 2023.",
            "[4] Song et al. \"RAG-HAT.\" EMNLP 2024.",
            "[5] Lin, Hilton, Evans. \"TruthfulQA.\" ACL 2022.",
            "[6] Lin. \"ROUGE.\" 2004.",
            "[7] Lewis et al. \"Retrieval-Augmented Generation.\" NeurIPS 2020.",
            "[8] Wei et al. \"Chain-of-Thought prompting.\" NeurIPS 2022.",
            "[9] AI@Meta. \"LLaMA 3.2 model card,\" 2024.",
            "[10] Ollama. \"Run LLMs locally,\" 2024.",
        ], body_size=12)

    # ---- Slide 12: Questions ---- (already says "Questions", just leave as-is)

    p.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
